#!/usr/bin/env python3
"""Generate work handover PPT for floxml_tools — focused on core workflow."""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
import os

# ── Colors ──
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DARK_BLUE = RGBColor(0x1B, 0x3A, 0x5C)
ACCENT_BLUE = RGBColor(0x2E, 0x75, 0xB6)
LIGHT_BLUE = RGBColor(0xD6, 0xE4, 0xF0)
MID_BLUE = RGBColor(0x5B, 0x9B, 0xD5)
DARK_GRAY = RGBColor(0x33, 0x33, 0x33)
MED_GRAY = RGBColor(0x66, 0x66, 0x66)
LIGHT_GRAY = RGBColor(0xF2, 0xF2, 0xF2)
ORANGE = RGBColor(0xED, 0x7D, 0x31)
GREEN = RGBColor(0x70, 0xAD, 0x47)
PURPLE = RGBColor(0x8E, 0x44, 0xAD)
RED = RGBColor(0xC0, 0x39, 0x2B)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


def _bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def _rect(slide, l, t, w, h, fill=None):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, w, h)
    s.line.fill.background()
    if fill:
        s.fill.solid()
        s.fill.fore_color.rgb = fill
    return s


def _txt(slide, l, t, w, h, text, sz=18, color=DARK_GRAY,
         bold=False, align=PP_ALIGN.LEFT, font="Microsoft YaHei"):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(sz)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font
    p.alignment = align
    return tb


def _bullets(slide, l, t, w, h, items, sz=14, color=DARK_GRAY, gap=Pt(6)):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.font.size = Pt(sz)
        p.font.color.rgb = color
        p.font.name = "Microsoft YaHei"
        p.space_after = gap
    return tb


def _footer(slide):
    _txt(slide, Inches(0.5), Inches(7.0), Inches(12), Inches(0.4),
         "floxml_tools 工作交接  |  2026-06-01", sz=9, color=MED_GRAY)


def _accent_bar(slide, color=ACCENT_BLUE):
    _rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.08), color)


def _section_title(slide, title, color=DARK_BLUE):
    _txt(slide, Inches(0.6), Inches(0.3), Inches(12), Inches(0.7),
         title, sz=32, color=color, bold=True)


# ========================================================================
# Slide 1: Cover
# ========================================================================
def slide_cover(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(s, WHITE)
    _accent_bar(s)
    _rect(s, Inches(1.5), Inches(1.8), Inches(10.3), Inches(3.6), LIGHT_BLUE)
    _txt(s, Inches(2), Inches(2.0), Inches(9.3), Inches(1.0),
         "floxml_tools", sz=48, color=DARK_BLUE, bold=True)
    _txt(s, Inches(2), Inches(3.0), Inches(9.3), Inches(0.8),
         "FloTHERM FloXML 自动化工具集 — 工作交接", sz=26, color=ACCENT_BLUE)
    _txt(s, Inches(2), Inches(4.0), Inches(9.3), Inches(0.6),
         "ECXML 导出 → FloXML 转换 → JSON 配置修改 → 求解", sz=18, color=MED_GRAY)
    _txt(s, Inches(2), Inches(5.5), Inches(9.3), Inches(0.4),
         "2026-06-01", sz=14, color=MED_GRAY)


# ========================================================================
# Slide 2: Workflow overview
# ========================================================================
def slide_workflow(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(s, WHITE)
    _accent_bar(s)
    _section_title(s, "核心工作流")

    _txt(s, Inches(0.6), Inches(1.1), Inches(12), Inches(0.5),
         "从项目导出 ECXML，转换为 FloXML，通过 JSON 配置修改后求解。",
         sz=16, color=MED_GRAY)

    # 4-step flow
    steps = [
        ("Step 1", "从 FloTHERM 项目\n导出 ECXML", "项目 → Export\n→ model.ecxml", ACCENT_BLUE),
        ("Step 2", "ECXML → FloXML\n格式转换", "ecxml_to_floxml\n_converter.py", MID_BLUE),
        ("Step 3", "JSON 配置\n修改 FloXML", "volume_regions\nsolve_settings\nboundary_conditions", GREEN),
        ("Step 4", "导入 FloTHERM\n求解", "完整项目文件\n→ Solve", RED),
    ]

    box_w = Inches(2.7)
    box_h = Inches(3.2)
    gap = Inches(0.3)
    sx = Inches(0.5)
    sy = Inches(1.8)

    for i, (label, title, detail, color) in enumerate(steps):
        x = sx + i * (box_w + gap)
        # Label
        _rect(s, x, sy, box_w, Inches(0.45), color)
        _txt(s, x, sy + Inches(0.05), box_w, Inches(0.35),
             label, sz=14, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
        # Body
        _rect(s, x, sy + Inches(0.45), box_w, box_h - Inches(0.45), LIGHT_GRAY)
        _txt(s, x + Inches(0.15), sy + Inches(0.55), box_w - Inches(0.3), Inches(1.0),
             title, sz=15, color=color, bold=True)
        _txt(s, x + Inches(0.15), sy + Inches(1.6), box_w - Inches(0.3), Inches(1.5),
             detail, sz=12, color=MED_GRAY)

        # Arrow
        if i < 3:
            ax = x + box_w + Inches(0.05)
            _txt(s, ax, sy + Inches(1.2), gap - Inches(0.1), Inches(0.5),
                 "→", sz=28, color=MED_GRAY, align=PP_ALIGN.CENTER)

    # JSON config summary at bottom
    _rect(s, Inches(0.5), Inches(5.3), Inches(12.1), Inches(1.5), LIGHT_BLUE)
    _txt(s, Inches(0.7), Inches(5.4), Inches(11.7), Inches(0.35),
         "Step 3 的三个 JSON 配置文件（与三个 Python 脚本一一对应）", sz=14, color=DARK_BLUE, bold=True)
    configs = [
        "• solve_settings.example.json  →  floxml_add_solve_settings.py  —  求解器参数、瞬态设置",
        "• floxml_volume_regions.example.json  →  floxml_add_volume_regions.py  —  体积区域、网格约束",
        "• boundary_conditions_*.json  →  floxml_boundary_conditions.py  —  环境条件、热源、表面属性、辐射",
    ]
    _bullets(s, Inches(0.7), Inches(5.8), Inches(11.7), Inches(1.0),
             configs, sz=12, color=DARK_GRAY, gap=Pt(3))

    _footer(s)


# ========================================================================
# Slide 3: ecxml_to_floxml_converter
# ========================================================================
def slide_ecxml_converter(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(s, WHITE)
    _accent_bar(s, MID_BLUE)

    _txt(s, Inches(0.6), Inches(0.3), Inches(9), Inches(0.5),
         "ecxml_to_floxml_converter.py", sz=28, color=DARK_BLUE, bold=True)
    _txt(s, Inches(0.6), Inches(0.8), Inches(9), Inches(0.4),
         "ECXML → FloXML 核心转换器  |  1,851 行", sz=16, color=MID_BLUE)
    _rect(s, Inches(0.6), Inches(1.2), Inches(12), Inches(0.02), MID_BLUE)

    _txt(s, Inches(0.6), Inches(1.4), Inches(12), Inches(0.4),
         "将 JEDEC JEP181 ECXML 器件热模型转换为 FloTHERM FloXML 项目格式，自动补充缺失配置。",
         sz=15, color=DARK_GRAY, bold=True)

    # Left: details
    _rect(s, Inches(0.5), Inches(2.0), Inches(7.2), Inches(4.8), LIGHT_GRAY)
    _txt(s, Inches(0.7), Inches(2.1), Inches(6.8), Inches(0.4),
         "功能说明", sz=16, color=DARK_BLUE, bold=True)
    details = [
        "• ECXML 是 JEDEC JEP181 器件级热模型交换标准格式",
        "• ECXML 缺少: 网格(grid)、求解器(solve)、模型设置(model)、求解域(solution_domain)",
        "• 本工具自动补充这些缺失部分，生成可直接导入 FloTHERM 的完整项目",
        "• 支持 JSON 配置文件自定义转换参数（网格、求解、区域等）",
        "• 完整保留 ECXML 中的几何体、材料、热源等定义",
    ]
    _bullets(s, Inches(0.7), Inches(2.5), Inches(6.8), Inches(2.0),
             details, sz=13, color=DARK_GRAY, gap=Pt(5))

    # Right: usage
    _rect(s, Inches(8.0), Inches(2.0), Inches(4.8), Inches(4.8), LIGHT_BLUE)
    _txt(s, Inches(8.2), Inches(2.1), Inches(4.4), Inches(0.4),
         "命令行用法", sz=16, color=DARK_BLUE, bold=True)
    usage = (
        "# 基本转换\n"
        "python -m floxml_tools.\n"
        "  ecxml_to_floxml_converter\n"
        "  model.ecxml -o project.xml\n\n"
        "# 带 JSON 配置\n"
        "python -m floxml_tools.\n"
        "  ecxml_to_floxml_converter\n"
        "  model.ecxml --config\n"
        "  config.json -o project.xml"
    )
    _txt(s, Inches(8.2), Inches(2.5), Inches(4.4), Inches(4.0),
         usage, sz=12, color=DARK_GRAY, font="Menlo")

    _footer(s)


# ========================================================================
# Slide 4: floxml_add_solve_settings
# ========================================================================
def slide_solve_settings(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(s, WHITE)
    _accent_bar(s, ORANGE)

    _txt(s, Inches(0.6), Inches(0.3), Inches(9), Inches(0.5),
         "floxml_add_solve_settings.py", sz=28, color=DARK_BLUE, bold=True)
    _txt(s, Inches(0.6), Inches(0.8), Inches(9), Inches(0.4),
         "求解设置注入  |  1,187 行", sz=16, color=ORANGE)
    _rect(s, Inches(0.6), Inches(1.2), Inches(12), Inches(0.02), ORANGE)

    _txt(s, Inches(0.6), Inches(1.4), Inches(12), Inches(0.4),
         "通过 JSON 或 Excel 配置向 FloXML 注入求解器设置和瞬态设置。",
         sz=15, color=DARK_GRAY, bold=True)

    # Left: what it configures
    _rect(s, Inches(0.5), Inches(2.0), Inches(7.2), Inches(4.8), LIGHT_GRAY)
    _txt(s, Inches(0.7), Inches(2.1), Inches(6.8), Inches(0.4),
         "可配置项", sz=16, color=DARK_BLUE, bold=True)

    items = [
        "• <model> 部分:",
        "    modeling — 求解类型 (flow_heat / conduction_only 等), 辐射开关",
        "    turbulence — 层流/湍流模型",
        "    gravity — 重力方向和大小",
        "    global — 基准压力、环境温度",
        "    initial_variables — 初始值设置",
        "    transient — 瞬态时间设置 (时间段、保存时刻、时间片)",
        "",
        "• <solve> 部分:",
        "    overall_control — 迭代次数、求解器类型、收敛判据",
        "    variable_controls — 各变量 (速度/温度) 的虚假时间步等",
        "    solver_controls — 压力/温度线性松弛因子",
    ]
    _bullets(s, Inches(0.7), Inches(2.5), Inches(6.8), Inches(4.0),
             items, sz=12, color=DARK_GRAY, gap=Pt(2))

    # Right: usage
    _rect(s, Inches(8.0), Inches(2.0), Inches(4.8), Inches(4.8), LIGHT_GRAY)
    _txt(s, Inches(8.2), Inches(2.1), Inches(4.4), Inches(0.4),
         "命令行用法", sz=16, color=DARK_BLUE, bold=True)
    usage = (
        "# JSON 配置注入\n"
        "python -m floxml_tools.\n"
        "  floxml_add_solve_settings\n"
        "  model.xml --config solve.json\n"
        "  -o output.xml\n\n"
        "# Excel 配置注入\n"
        "python -m floxml_tools.\n"
        "  floxml_add_solve_settings\n"
        "  model.xml --config solve.xlsx\n"
        "  -o output.xml\n\n"
        "# 生成配置模板\n"
        "python -m floxml_tools.\n"
        "  floxml_add_solve_settings\n"
        "  --create-template template.xlsx"
    )
    _txt(s, Inches(8.2), Inches(2.5), Inches(4.4), Inches(4.0),
         usage, sz=11, color=DARK_GRAY, font="Menlo")

    _footer(s)


# ========================================================================
# Slide 5: solve_settings JSON config detail
# ========================================================================
def slide_solve_config(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(s, WHITE)
    _accent_bar(s, ORANGE)
    _section_title(s, "solve_settings.example.json 配置详解")

    # Left half: model section
    _rect(s, Inches(0.5), Inches(1.1), Inches(6.0), Inches(5.8), LIGHT_GRAY)
    _txt(s, Inches(0.7), Inches(1.2), Inches(5.6), Inches(0.35),
         "model 部分 — 求解模型设置", sz=14, color=ORANGE, bold=True)
    model_json = (
        '"modeling": {\n'
        '  "solution": "flow_heat",   // flow_heat / conduction_only\n'
        '  "radiation": "on",          // on / off\n'
        '  "dimensionality": "3d"\n'
        '},\n'
        '"turbulence": { "type": "laminar" },\n'
        '"gravity": {\n'
        '  "type": "normal",\n'
        '  "normal_direction": "neg_y",\n'
        '  "value": 9.81\n'
        '},\n'
        '"global": {\n'
        '  "datum_pressure": 101325.0,    // Pa\n'
        '  "ambient_temperature": 25.0     // °C\n'
        '}'
    )
    _txt(s, Inches(0.7), Inches(1.6), Inches(5.6), Inches(5.0),
         model_json, sz=11, color=DARK_GRAY, font="Menlo")

    # Right half: solve section
    _rect(s, Inches(6.8), Inches(1.1), Inches(6.0), Inches(5.8), LIGHT_GRAY)
    _txt(s, Inches(7.0), Inches(1.2), Inches(5.6), Inches(0.35),
         "solve 部分 — 迭代控制 + 瞬态", sz=14, color=ORANGE, bold=True)
    solve_json = (
        '"overall_control": {\n'
        '  "outer_iterations": 500,\n'
        '  "solver_option": "multi_grid",\n'
        '  "convergence_values": {\n'
        '    "required_accuracy": 0.2,\n'
        '    "num_iterations": 45\n'
        '  }\n'
        '},\n'
        '"transient": {\n'
        '  "overall_transient": {\n'
        '    "start_time": 0,\n'
        '    "end_time": 60\n'
        '  },\n'
        '  "save_times": [0, 30, 60],\n'
        '  "time_patches": [{\n'
        '    "name": "First",\n'
        '    "start_time": 0,\n'
        '    "end_time": 30,\n'
        '    "step_control": "minimum_number",\n'
        '    "minimum_number": 15\n'
        '  }]\n'
        '}'
    )
    _txt(s, Inches(7.0), Inches(1.6), Inches(5.6), Inches(5.0),
         solve_json, sz=11, color=DARK_GRAY, font="Menlo")

    _footer(s)


# ========================================================================
# Slide 6: floxml_add_volume_regions
# ========================================================================
def slide_volume_regions(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(s, WHITE)
    _accent_bar(s, GREEN)

    _txt(s, Inches(0.6), Inches(0.3), Inches(9), Inches(0.5),
         "floxml_add_volume_regions.py", sz=28, color=DARK_BLUE, bold=True)
    _txt(s, Inches(0.6), Inches(0.8), Inches(9), Inches(0.4),
         "体积区域 & 网格约束注入  |  1,225 行", sz=16, color=GREEN)
    _rect(s, Inches(0.6), Inches(1.2), Inches(12), Inches(0.02), GREEN)

    _txt(s, Inches(0.6), Inches(1.4), Inches(12), Inches(0.4),
         "通过 JSON 配置向 FloXML 添加体积区域 (region) 和网格约束 (grid_constraint)。",
         sz=15, color=DARK_GRAY, bold=True)

    # Left: features
    _rect(s, Inches(0.5), Inches(2.0), Inches(7.2), Inches(4.8), LIGHT_GRAY)
    _txt(s, Inches(0.7), Inches(2.1), Inches(6.8), Inches(0.4),
         "两种区域定义方式", sz=16, color=DARK_BLUE, bold=True)

    items = [
        "1. 显式定义 — 直接指定位置和尺寸:",
        "     position: [x, y, z]    size: [sx, sy, sz]",
        "",
        "2. 从几何体包围盒自动推导 (bbox_from):",
        "     include_names: [\"PCB\", \"U1\"]   — 按名称匹配",
        "     include_patterns: [\"R22*\", \"C*\"] — 通配符匹配",
        "     include_tags: [\"cuboid\"]         — 按类型匹配",
        "     padding: [px, py, pz]             — 外扩余量",
        "",
        "网格约束 (grid_constraint):",
        "• 先在 grid_constraints 中定义约束（最小尺寸、最小数量等）",
        "• 再在 regions 中引用约束名，应用到区域或几何体",
        "• object_constraints 可直接给指定几何体分配约束",
        "",
        "区域可放置在根 geometry 或指定 assembly 下",
    ]
    _bullets(s, Inches(0.7), Inches(2.5), Inches(6.8), Inches(4.0),
             items, sz=12, color=DARK_GRAY, gap=Pt(2))

    # Right: usage
    _rect(s, Inches(8.0), Inches(2.0), Inches(4.8), Inches(4.8), LIGHT_GRAY)
    _txt(s, Inches(8.2), Inches(2.1), Inches(4.4), Inches(0.4),
         "命令行用法", sz=16, color=DARK_BLUE, bold=True)
    usage = (
        "# 注入体积区域和网格约束\n"
        "python -m floxml_tools.\n"
        "  floxml_add_volume_regions\n"
        "  model.xml --config\n"
        "  regions.json -o output.xml\n\n"
        "# 原地修改\n"
        "python -m floxml_tools.\n"
        "  floxml_add_volume_regions\n"
        "  model.xml --config\n"
        "  regions.json --in-place"
    )
    _txt(s, Inches(8.2), Inches(2.5), Inches(4.4), Inches(4.0),
         usage, sz=11, color=DARK_GRAY, font="Menlo")

    _footer(s)


# ========================================================================
# Slide 7: volume_regions JSON config detail
# ========================================================================
def slide_volume_config(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(s, WHITE)
    _accent_bar(s, GREEN)
    _section_title(s, "floxml_volume_regions.example.json 配置详解")

    # Left: grid constraints + object constraints
    _rect(s, Inches(0.5), Inches(1.1), Inches(6.0), Inches(5.8), LIGHT_GRAY)
    _txt(s, Inches(0.7), Inches(1.2), Inches(5.6), Inches(0.35),
         "网格约束定义 + 几何体约束分配", sz=14, color=GREEN, bold=True)
    left_json = (
        '"grid_constraints": [\n'
        '  {\n'
        '    "name": "Grid Constraint 1",\n'
        '    "enable_min_cell_size": true,\n'
        '    "min_cell_size": 0.001,\n'
        '    "min_number": 43,\n'
        '    "high_inflation": {\n'
        '      "inflation_type": "size",\n'
        '      "inflation_size": 0.005,\n'
        '      "min_number": 23\n'
        '    }\n'
        '  }\n'
        '],\n'
        '"object_constraints": [\n'
        '  {\n'
        '    "target_names": ["PCB"],\n'
        '    "all_grid_constraint":\n'
        '      "Grid Constraint 1"\n'
        '  }\n'
        ']'
    )
    _txt(s, Inches(0.7), Inches(1.6), Inches(5.6), Inches(5.0),
         left_json, sz=11, color=DARK_GRAY, font="Menlo")

    # Right: regions
    _rect(s, Inches(6.8), Inches(1.1), Inches(6.0), Inches(5.8), LIGHT_GRAY)
    _txt(s, Inches(7.0), Inches(1.2), Inches(5.6), Inches(0.35),
         "体积区域定义 (显式 + bbox)", sz=14, color=GREEN, bold=True)
    right_json = (
        '"regions": [\n'
        '  {\n'
        '    "name": "Explicit Volume Region",\n'
        '    "position": [-0.01, -0.01, -0.002],\n'
        '    "size": [0.12, 0.08, 0.01],\n'
        '    "localized_grid": true,\n'
        '    "x_grid_constraint":\n'
        '      "Grid Constraint 1"\n'
        '  },\n'
        '  {\n'
        '    "name": "BBox Region Around PCB",\n'
        '    "parent_assembly": "DemoBoard",\n'
        '    "bbox_from": {\n'
        '      "include_names": ["PCB"],\n'
        '      "include_patterns": ["U*"],\n'
        '      "padding": [0.001, 0.001, 0.0005]\n'
        '    },\n'
        '    "all_grid_constraint":\n'
        '      "Grid Constraint 1"\n'
        '  }\n'
        ']'
    )
    _txt(s, Inches(7.0), Inches(1.6), Inches(5.6), Inches(5.0),
         right_json, sz=11, color=DARK_GRAY, font="Menlo")

    _footer(s)


# ========================================================================
# Slide 8: floxml_boundary_conditions
# ========================================================================
def slide_boundary_conditions(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(s, WHITE)
    _accent_bar(s, PURPLE)

    _txt(s, Inches(0.6), Inches(0.3), Inches(9), Inches(0.5),
         "floxml_boundary_conditions.py", sz=28, color=DARK_BLUE, bold=True)
    _txt(s, Inches(0.6), Inches(0.8), Inches(9), Inches(0.4),
         "边界条件注入  |  921 行", sz=16, color=PURPLE)
    _rect(s, Inches(0.6), Inches(1.2), Inches(12), Inches(0.02), PURPLE)

    _txt(s, Inches(0.6), Inches(1.4), Inches(12), Inches(0.4),
         "通过 JSON 配置向 FloXML 添加/修改 7 种边界条件类型。",
         sz=15, color=DARK_GRAY, bold=True)

    # 7 type cards
    types = [
        ("Ambient\n(环境)", "温度、压力、换热系数\n速度、辐射温度", ACCENT_BLUE),
        ("Solution Domain\n(求解域)", "各面边界类型\nambient/symmetry/\nwall/opening", MID_BLUE),
        ("Surface\n(表面属性)", "发射率、粗糙度\narea_factor", GREEN),
        ("Radiation\n(辐射)", "表面类型\n面积阈值", ORANGE),
        ("Source\n(热源)", "总功率 / 定温\n体积热源", RED),
        ("Surface Exchange\n(表面交换)", "对流换热方法\nHTC 常数/曲线", PURPLE),
        ("Thermal\n(热模型)", "导热/对流模型", RGBColor(0x2C, 0x3E, 0x50)),
    ]

    cw = Inches(1.65)
    ch = Inches(2.8)
    cg = Inches(0.12)
    cx = Inches(0.4)
    cy = Inches(2.0)

    for i, (name, desc, color) in enumerate(types):
        x = cx + i * (cw + cg)
        _rect(s, x, cy, cw, Inches(0.6), color)
        _txt(s, x + Inches(0.05), cy + Inches(0.05), cw - Inches(0.1), Inches(0.5),
             name, sz=11, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
        _rect(s, x, cy + Inches(0.6), cw, ch - Inches(0.6), LIGHT_GRAY)
        _txt(s, x + Inches(0.08), cy + Inches(0.7), cw - Inches(0.16), ch - Inches(0.8),
             desc, sz=10, color=DARK_GRAY)

    # Bottom: face_conditions + usage
    _rect(s, Inches(0.5), Inches(5.0), Inches(7.0), Inches(2.0), LIGHT_BLUE)
    _txt(s, Inches(0.7), Inches(5.1), Inches(6.6), Inches(0.35),
         "face_conditions — 将属性分配到几何体", sz=13, color=DARK_BLUE, bold=True)
    face_desc = [
        "• target_names / target_patterns / target_tags 匹配几何体",
        "• 将 surface / radiation / thermal / source 等属性分配上去",
        "• 例: Chip_U1 → Paint 表面 + Conduction 热模型",
    ]
    _bullets(s, Inches(0.7), Inches(5.5), Inches(6.6), Inches(1.3),
             face_desc, sz=11, color=DARK_GRAY, gap=Pt(2))

    _rect(s, Inches(7.8), Inches(5.0), Inches(5.0), Inches(2.0), LIGHT_GRAY)
    _txt(s, Inches(8.0), Inches(5.1), Inches(4.6), Inches(0.35),
         "命令行", sz=13, color=DARK_BLUE, bold=True)
    usage = (
        "python -m floxml_tools.\n"
        "  floxml_boundary_conditions\n"
        "  model.xml --config bc.json\n"
        "  -o output.xml"
    )
    _txt(s, Inches(8.0), Inches(5.5), Inches(4.6), Inches(1.3),
         usage, sz=11, color=DARK_GRAY, font="Menlo")

    _footer(s)


# ========================================================================
# Slide 9: boundary_conditions JSON config detail
# ========================================================================
def slide_boundary_config(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(s, WHITE)
    _accent_bar(s, PURPLE)
    _section_title(s, "boundary_conditions JSON 配置详解")

    # Left: ambients + solution_domain
    _rect(s, Inches(0.5), Inches(1.1), Inches(6.0), Inches(5.8), LIGHT_GRAY)
    _txt(s, Inches(0.7), Inches(1.2), Inches(5.6), Inches(0.35),
         "环境条件 + 求解域边界", sz=14, color=PURPLE, bold=True)
    left_json = (
        '"ambients": [\n'
        '  {\n'
        '    "name": "Ambient",\n'
        '    "pressure": 101325,\n'
        '    "temperature": 293,      // K\n'
        '    "heat_transfer_coeff": 12,\n'
        '    "velocity": [0, 0, 0]\n'
        '  },\n'
        '  {\n'
        '    "name": "Forced Air",\n'
        '    "temperature": 300,\n'
        '    "velocity": [0, 0, 2.5],  // m/s\n'
        '    "heat_transfer_coeff": 25\n'
        '  }\n'
        '],\n'
        '"solution_domain": {\n'
        '  "x_low_boundary": "ambient",\n'
        '  "y_low_boundary": "symmetry",\n'
        '  "z_high_boundary": "opening",\n'
        '  "x_low_ambient": "Ambient",\n'
        '  "z_high_ambient": "Forced Air"\n'
        '}'
    )
    _txt(s, Inches(0.7), Inches(1.6), Inches(5.6), Inches(5.0),
         left_json, sz=11, color=DARK_GRAY, font="Menlo")

    # Right: sources + face_conditions
    _rect(s, Inches(6.8), Inches(1.1), Inches(6.0), Inches(2.6), LIGHT_GRAY)
    _txt(s, Inches(7.0), Inches(1.2), Inches(5.6), Inches(0.35),
         "热源定义", sz=14, color=PURPLE, bold=True)
    source_json = (
        '"sources": [\n'
        '  {\n'
        '    "name": "Heat Source",\n'
        '    "source_options": [{\n'
        '      "applies_to": "temperature",\n'
        '      "type": "total",\n'
        '      "value": 23.3       // W\n'
        '    }]\n'
        '  }\n'
        ']'
    )
    _txt(s, Inches(7.0), Inches(1.6), Inches(5.6), Inches(2.0),
         source_json, sz=11, color=DARK_GRAY, font="Menlo")

    # Bottom right: face_conditions
    _rect(s, Inches(6.8), Inches(3.9), Inches(6.0), Inches(3.0), LIGHT_BLUE)
    _txt(s, Inches(7.0), Inches(4.0), Inches(5.6), Inches(0.35),
         "属性分配 (face_conditions)", sz=14, color=PURPLE, bold=True)
    face_json = (
        '"face_conditions": [\n'
        '  {\n'
        '    "target_names": ["Heatsink"],\n'
        '    "surface": "Polished Metal",\n'
        '    "radiation": "Sub-Divided",\n'
        '    "thermal": "Conduction"\n'
        '  },\n'
        '  {\n'
        '    "target_names": ["Chip_U1"],\n'
        '    "surface": "Paint",\n'
        '    "surface_exchange": "Constant HTC"\n'
        '  }\n'
        ']'
    )
    _txt(s, Inches(7.0), Inches(4.4), Inches(5.6), Inches(2.3),
         face_json, sz=11, color=DARK_GRAY, font="Menlo")

    _footer(s)


# ========================================================================
# Slide 10: Typical usage — complete workflow
# ========================================================================
def slide_typical_usage(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(s, WHITE)
    _accent_bar(s)
    _section_title(s, "典型使用流程")

    steps = [
        ("1. 从 FloTHERM 导出 ECXML",
         "在 FloTHERM 中打开项目 → File → Export → ECXML\n保存为 model.ecxml",
         ACCENT_BLUE),
        ("2. ECXML → FloXML 转换",
         "python -m floxml_tools.ecxml_to_floxml_converter \\\n"
         "  model.ecxml -o project.xml",
         MID_BLUE),
        ("3. 注入求解设置",
         "python -m floxml_tools.floxml_add_solve_settings \\\n"
         "  project.xml --config solve_settings.json -o project.xml",
         ORANGE),
        ("4. 注入体积区域和网格约束",
         "python -m floxml_tools.floxml_add_volume_regions \\\n"
         "  project.xml --config regions.json --in-place",
         GREEN),
        ("5. 注入边界条件",
         "python -m floxml_tools.floxml_boundary_conditions \\\n"
         "  project.xml --config boundary.json --in-place",
         PURPLE),
        ("6. 导入 FloTHERM 求解",
         "在 FloTHERM 中 File → Import → 选择最终 project.xml → Solve",
         RED),
    ]

    sy = Inches(1.1)
    step_h = Inches(0.95)
    for i, (title, code, color) in enumerate(steps):
        y = sy + i * step_h
        # Step label
        _rect(s, Inches(0.4), y, Inches(3.0), step_h - Inches(0.08), color)
        _txt(s, Inches(0.55), y + Inches(0.12), Inches(2.7), Inches(0.6),
             title, sz=13, color=WHITE, bold=True)
        # Code
        _rect(s, Inches(3.6), y, Inches(9.3), step_h - Inches(0.08), LIGHT_GRAY)
        _txt(s, Inches(3.8), y + Inches(0.08), Inches(8.9), Inches(0.7),
             code, sz=11, color=DARK_GRAY, font="Menlo")

    _footer(s)


# ========================================================================
# Slide 11: floxml_pipeline — one-pass pipeline
# ========================================================================
def slide_pipeline(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(s, WHITE)
    _accent_bar(s, RGBColor(0x2C, 0x3E, 0x50))

    _txt(s, Inches(0.6), Inches(0.3), Inches(9), Inches(0.5),
         "floxml_pipeline.py", sz=28, color=DARK_BLUE, bold=True)
    _txt(s, Inches(0.6), Inches(0.8), Inches(9), Inches(0.4),
         "一键流水线 — 三个修改脚本合并为一次执行  |  308 行", sz=16, color=RGBColor(0x2C, 0x3E, 0x50))
    _rect(s, Inches(0.6), Inches(1.2), Inches(12), Inches(0.02), RGBColor(0x2C, 0x3E, 0x50))

    _txt(s, Inches(0.6), Inches(1.4), Inches(12), Inches(0.4),
         "将前面 3 个注入步骤合并成一条命令，支持分步或统一配置。",
         sz=15, color=DARK_GRAY, bold=True)

    # Left: execution flow
    _rect(s, Inches(0.5), Inches(2.0), Inches(7.2), Inches(4.8), LIGHT_GRAY)
    _txt(s, Inches(0.7), Inches(2.1), Inches(6.8), Inches(0.35),
         "执行流程", sz=16, color=DARK_BLUE, bold=True)

    flow_items = [
        ("Step 1", "--wrap", "包装几何文件为完整项目 (可选)", ACCENT_BLUE),
        ("Step 2", "--grid", "注入体积区域 + 网格约束", GREEN),
        ("Step 3", "--solve", "注入求解设置 + 瞬态设置", ORANGE),
        ("Step 4", "--boundary", "注入边界条件 (环境/热源/表面/辐射)", PURPLE),
    ]
    for i, (label, flag, desc, color) in enumerate(flow_items):
        y = Inches(2.6) + i * Inches(0.9)
        _rect(s, Inches(0.7), y, Inches(1.0), Inches(0.5), color)
        _txt(s, Inches(0.7), y + Inches(0.07), Inches(1.0), Inches(0.35),
             label, sz=11, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
        _txt(s, Inches(1.9), y + Inches(0.07), Inches(1.2), Inches(0.35),
             flag, sz=12, color=color, bold=True, font="Menlo")
        _txt(s, Inches(3.2), y + Inches(0.07), Inches(4.2), Inches(0.35),
             desc, sz=12, color=DARK_GRAY)
        if i < 3:
            _txt(s, Inches(1.1), y + Inches(0.5), Inches(0.5), Inches(0.3),
                 "↓", sz=16, color=MED_GRAY, align=PP_ALIGN.CENTER)

    # unified config note
    _rect(s, Inches(0.7), Inches(6.1), Inches(6.8), Inches(0.5), LIGHT_BLUE)
    _txt(s, Inches(0.9), Inches(6.15), Inches(6.4), Inches(0.4),
         "-c config.json 会按 key 自动拆分到 grid / solve / boundary",
         sz=11, color=DARK_BLUE, bold=True)

    # Right: usage examples
    _rect(s, Inches(8.0), Inches(2.0), Inches(4.8), Inches(4.8), LIGHT_GRAY)
    _txt(s, Inches(8.2), Inches(2.1), Inches(4.4), Inches(0.35),
         "命令行用法", sz=16, color=DARK_BLUE, bold=True)

    usage = (
        "# 统一配置一步到位\n"
        "python -m floxml_tools.\n"
        "  floxml_pipeline project.xml \\\n"
        "  -c config.json -o output.xml\n\n"
        "# 分别指定三个配置\n"
        "python -m floxml_tools.\n"
        "  floxml_pipeline project.xml \\\n"
        "  --grid regions.json \\\n"
        "  --solve solve.json \\\n"
        "  --boundary bc.json \\\n"
        "  -o output.xml\n\n"
        "# 只跑其中一步\n"
        "python -m floxml_tools.\n"
        "  floxml_pipeline project.xml \\\n"
        "  --boundary bc.json \\\n"
        "  -o output.xml"
    )
    _txt(s, Inches(8.2), Inches(2.5), Inches(4.4), Inches(4.0),
         usage, sz=10, color=DARK_GRAY, font="Menlo")

    _footer(s)


# ========================================================================
# Slide 12: JSON config files summary
# ========================================================================
def slide_config_summary(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(s, WHITE)
    _accent_bar(s)
    _section_title(s, "JSON 配置文件汇总")

    _txt(s, Inches(0.6), Inches(1.1), Inches(12), Inches(0.4),
         "每个 Python 脚本对应一个或多个 JSON 配置文件，修改 JSON 即可调整 FloXML 项目参数，无需手动编辑 XML。",
         sz=14, color=MED_GRAY)

    # Table-like layout
    headers = ["配置文件", "对应脚本", "用途"]
    col_widths = [Inches(3.8), Inches(3.8), Inches(4.4)]
    col_x = [Inches(0.5), Inches(4.4), Inches(8.3)]

    # Header row
    _rect(s, Inches(0.5), Inches(1.7), Inches(12.1), Inches(0.5), ACCENT_BLUE)
    for j, h in enumerate(headers):
        _txt(s, col_x[j], Inches(1.75), col_widths[j], Inches(0.4),
             h, sz=14, color=WHITE, bold=True)

    rows = [
        ("solve_settings.example.json", "floxml_add_solve_settings.py", "求解器参数、收敛判据、瞬态时间"),
        ("solve_settings.example.xlsx", "floxml_add_solve_settings.py", "同上，Excel 格式（可用 --create-template 生成）"),
        ("transient_only.example.json", "floxml_add_solve_settings.py", "仅瞬态设置（start/end_time, save_times, patches）"),
        ("transient_maxsize.example.json", "floxml_add_solve_settings.py", "瞬态设置（max_size 控制步长）"),
        ("floxml_volume_regions.example.json", "floxml_add_volume_regions.py", "体积区域定义 + 网格约束"),
        ("boundary_conditions_*.json", "floxml_boundary_conditions.py", "环境条件、热源、表面属性、辐射、表面交换"),
    ]

    for i, (cfg, script, desc) in enumerate(rows):
        y = Inches(2.3) + i * Inches(0.6)
        bg = LIGHT_GRAY if i % 2 == 0 else WHITE
        _rect(s, Inches(0.5), y, Inches(12.1), Inches(0.55), bg)
        _txt(s, col_x[0], y + Inches(0.08), col_widths[0], Inches(0.4),
             cfg, sz=12, color=ACCENT_BLUE, font="Menlo")
        _txt(s, col_x[1], y + Inches(0.08), col_widths[1], Inches(0.4),
             script, sz=12, color=DARK_GRAY, font="Menlo")
        _txt(s, col_x[2], y + Inches(0.08), col_widths[2], Inches(0.4),
             desc, sz=12, color=DARK_GRAY)

    _footer(s)


# ========================================================================
# Slide 12: XSD Schema — overview
# ========================================================================
def slide_xsd_overview(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(s, WHITE)
    _accent_bar(s, RGBColor(0x34, 0x49, 0x5E))
    _section_title(s, "FloXML XSD Schema 参考文档")

    _txt(s, Inches(0.6), Inches(1.1), Inches(12), Inches(0.4),
         "examples/DCIM Development Toolkit/Schema Files/FloXML/ — 4,001 行 XSD 定义了 FloXML 完整结构规范。",
         sz=14, color=MED_GRAY)

    # File overview cards
    xsd_files = [
        ("xmlSchema.xsd", "140 行", "主入口", "定义 <xml_case> 根元素\nmodel / solve / grid /\nattributes / geometry /\nsolution_domain 的\n嵌套顺序", ACCENT_BLUE),
        ("XmlDefinitions.xsd", "331 行", "公共类型", "基础类型库:\ntriplet / direction /\ntrueFalse / ratio /\nvariable_types 等\n枚举值定义", MID_BLUE),
        ("XmlAttributes.xsd", "783 行", "属性定义", "材料 / 热源 / 环境 /\n表面 / 辐射 / 风扇 /\n流阻 / 网格约束 /\n瞬态函数\n每种属性的字段和枚举", GREEN),
        ("XmlEntities.xsd", "615 行", "求解/网格", "model 求解模式\nsolve 求解器控制\ngrid 系统网格\npatches 局部加密\n全部字段和合法值", ORANGE),
        ("XmlGeometry.xsd", "2,132 行", "几何体", "cuboid / plate / prism /\nassembly / fan / vent /\nsource / monitor_point\n每个几何体的字段\n和可引用的属性", PURPLE),
    ]

    cw = Inches(2.3)
    ch = Inches(4.4)
    cg = Inches(0.12)
    cx = Inches(0.4)
    cy = Inches(1.7)

    for i, (fname, lines, role, desc, color) in enumerate(xsd_files):
        x = cx + i * (cw + cg)
        _rect(s, x, cy, cw, Inches(0.45), color)
        _txt(s, x + Inches(0.05), cy + Inches(0.05), cw - Inches(0.1), Inches(0.35),
             f"{role}  ({lines})", sz=11, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
        _rect(s, x, cy + Inches(0.45), cw, ch - Inches(0.45), LIGHT_GRAY)
        _txt(s, x + Inches(0.08), cy + Inches(0.55), cw - Inches(0.16), Inches(0.4),
             fname, sz=10, color=color, bold=True, font="Menlo")
        _txt(s, x + Inches(0.08), cy + Inches(1.0), cw - Inches(0.16), ch - Inches(1.1),
             desc, sz=10, color=DARK_GRAY)

    # Bottom: reading order + XSD syntax tips
    _rect(s, Inches(0.5), Inches(6.3), Inches(5.8), Inches(0.6), LIGHT_BLUE)
    _txt(s, Inches(0.7), Inches(6.35), Inches(5.4), Inches(0.5),
         "推荐阅读顺序: Definitions → Attributes → Entities → Geometry → Schema",
         sz=12, color=DARK_BLUE, bold=True)

    _rect(s, Inches(6.6), Inches(6.3), Inches(6.0), Inches(0.6), LIGHT_GRAY)
    _txt(s, Inches(6.8), Inches(6.35), Inches(5.6), Inches(0.5),
         "用途: 写 JSON 配置时查 XSD 确认字段名、枚举值、必填/可选",
         sz=12, color=DARK_GRAY)

    _footer(s)


# ========================================================================
# Slide 13: XSD Schema — key details
# ========================================================================
def slide_xsd_details(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(s, WHITE)
    _accent_bar(s, RGBColor(0x34, 0x49, 0x5E))
    _section_title(s, "XSD Schema — 关键内容速查")

    # Left: xml_case structure
    _rect(s, Inches(0.5), Inches(1.1), Inches(4.0), Inches(5.8), LIGHT_GRAY)
    _txt(s, Inches(0.7), Inches(1.2), Inches(3.6), Inches(0.35),
         "xml_case 整体结构", sz=14, color=ACCENT_BLUE, bold=True)
    structure = (
        "<xml_case>\n"
        "  <name/>             必填\n"
        "  <model/>            求解模式\n"
        "  <solve/>            求解器\n"
        "  <grid/>             网格\n"
        "  <attributes>\n"
        "    materials         材料库\n"
        "    surfaces          表面属性\n"
        "    surface_exchanges 表面换热\n"
        "    thermals          热模型\n"
        "    sources           热源\n"
        "    ambients          环境条件\n"
        "    fluids            流体\n"
        "    grid_constraints  网格约束\n"
        "    radiations        辐射模型\n"
        "    transients        瞬态函数\n"
        "    fans / resistances\n"
        "  </attributes>\n"
        "  <geometry/>         几何模型\n"
        "  <solution_domain/>  求解域\n"
        "</xml_case>"
    )
    _txt(s, Inches(0.7), Inches(1.6), Inches(3.6), Inches(5.0),
         structure, sz=10, color=DARK_GRAY, font="Menlo")

    # Middle: attribute types
    _rect(s, Inches(4.7), Inches(1.1), Inches(4.0), Inches(5.8), LIGHT_GRAY)
    _txt(s, Inches(4.9), Inches(1.2), Inches(3.6), Inches(0.35),
         "属性类型 (Attributes)", sz=14, color=GREEN, bold=True)
    attrs = (
        "isotropic_material_att\n"
        "  各向同性材料 (k, rho, cp)\n"
        "orthotropic_material_att\n"
        "  正交各向异性 (kx, ky, kz)\n"
        "biaxial_material_att\n"
        "  双轴材料 (in_plane, normal)\n"
        "\n"
        "source_att\n"
        "  total / volume / area / fixed\n"
        "  / linear / non_linear\n"
        "\n"
        "surface_att\n"
        "  emissivity, roughness\n"
        "radiation_att\n"
        "  non_radiating / single /\n"
        "  subdivided_radiating\n"
        "ambient_att\n"
        "  temperature, pressure,\n"
        "  velocity, HTC\n"
        "thermal_att\n"
        "  conduction / convection /\n"
        "  fixed_temperature\n"
        "fan_att\n"
        "  normal / angled / swirl\n"
        "  + fan_curve_points"
    )
    _txt(s, Inches(4.9), Inches(1.6), Inches(3.6), Inches(5.0),
         attrs, sz=10, color=DARK_GRAY, font="Menlo")

    # Right: geometry + XSD reading tips
    _rect(s, Inches(8.9), Inches(1.1), Inches(4.0), Inches(3.4), LIGHT_GRAY)
    _txt(s, Inches(9.1), Inches(1.2), Inches(3.6), Inches(0.35),
         "几何体类型", sz=14, color=PURPLE, bold=True)
    geom = (
        "cuboid       长方体 (最常用)\n"
        "plate        薄板 (PCB 等)\n"
        "prism        棱柱体\n"
        "assembly     装配体 (嵌套)\n"
        "fan          风扇\n"
        "vent         通风口\n"
        "source       热源几何\n"
        "monitor_point 监控点\n"
        "sloping_block 斜块\n"
        "opening      开口\n"
        "\n"
        "每个几何体可引用:\n"
        " material / thermal / surface\n"
        " radiation / source / fan\n"
        " ambient / resistance\n"
        " grid_constraint (按面)"
    )
    _txt(s, Inches(9.1), Inches(1.6), Inches(3.6), Inches(2.8),
         geom, sz=10, color=DARK_GRAY, font="Menlo")

    # XSD syntax tips
    _rect(s, Inches(8.9), Inches(4.7), Inches(4.0), Inches(2.2), LIGHT_BLUE)
    _txt(s, Inches(9.1), Inches(4.8), Inches(3.6), Inches(0.35),
         "XSD 语法速查", sz=14, color=DARK_BLUE, bold=True)
    tips = (
        'minOccurs="0"    可选字段\n'
        'minOccurs="1"    必填字段\n'
        'enumeration      合法枚举值\n'
        'restriction      值约束 (范围)\n'
        '<xs:all>         无序, 最多一次\n'
        '<xs:sequence>    有序\n'
        '<xs:choice>      多选一\n'
        'maxOccurs="unbounded"  可重复'
    )
    _txt(s, Inches(9.1), Inches(5.2), Inches(3.6), Inches(1.5),
         tips, sz=10, color=DARK_GRAY, font="Menlo")

    _footer(s)


# ========================================================================
# Slide 14: Command-line solving
# ========================================================================
def slide_cli_solve(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(s, WHITE)
    _accent_bar(s, RED)
    _section_title(s, "命令行求解 & 结果保存")

    _txt(s, Inches(0.6), Inches(1.1), Inches(12), Inches(0.4),
         "修改完 FloXML 后，通过命令行调用 FloTHERM 求解，保存为 .pack 项目和 HTML 求解报告。",
         sz=14, color=MED_GRAY)

    # Single file solve
    _rect(s, Inches(0.5), Inches(1.7), Inches(6.0), Inches(2.3), LIGHT_GRAY)
    _txt(s, Inches(0.7), Inches(1.8), Inches(5.6), Inches(0.35),
         "单文件求解", sz=15, color=RED, bold=True)
    solve_cmd = (
        "# ECXML 直接求解 → 保存 .pack + HTML 报告\n"
        "flotherm -b model.ecxml \\\n"
        "  -z output.pack -r report.html\n\n"
        "# FloXML 求解\n"
        "flotherm -b project.xml\n\n"
        "# FloSCRIPT 脚本求解 (最灵活)\n"
        "flotherm -b -f solve_script.xml"
    )
    _txt(s, Inches(0.7), Inches(2.2), Inches(5.6), Inches(1.6),
         solve_cmd, sz=11, color=DARK_GRAY, font="Menlo")

    # Key flags
    _rect(s, Inches(6.8), Inches(1.7), Inches(6.0), Inches(2.3), LIGHT_BLUE)
    _txt(s, Inches(7.0), Inches(1.8), Inches(5.6), Inches(0.35),
         "关键参数说明", sz=15, color=DARK_BLUE, bold=True)
    flags = [
        "-b              batch mode (无 GUI)",
        "-f script.xml   执行 FloSCRIPT 脚本",
        "-z output.pack  保存结果为 .pack 文件",
        "-r report.html  生成 HTML 求解报告",
        "--timeout N     求解超时 (秒)",
    ]
    _bullets(s, Inches(7.0), Inches(2.2), Inches(5.6), Inches(1.6),
             flags, sz=12, color=DARK_GRAY, gap=Pt(4))

    # Pack file operations
    _rect(s, Inches(0.5), Inches(4.2), Inches(6.0), Inches(2.8), LIGHT_GRAY)
    _txt(s, Inches(0.7), Inches(4.3), Inches(5.6), Inches(0.35),
         ".pack 文件操作", sz=15, color=ACCENT_BLUE, bold=True)
    pack_cmd = (
        "# 查看 pack 内容\n"
        "python pack_editor.py model.pack --list\n\n"
        "# 解压 pack 文件\n"
        "python pack_editor.py model.pack \\\n"
        "  --extract ./extracted\n\n"
        "# 修改 pack 中的功耗\n"
        "python pack_editor.py model.pack \\\n"
        "  --set-power U1_CPU 15.0 \\\n"
        "  -o modified.pack"
    )
    _txt(s, Inches(0.7), Inches(4.7), Inches(5.6), Inches(2.0),
         pack_cmd, sz=11, color=DARK_GRAY, font="Menlo")

    # Python solve tools
    _rect(s, Inches(6.8), Inches(4.2), Inches(6.0), Inches(2.8), LIGHT_GRAY)
    _txt(s, Inches(7.0), Inches(4.3), Inches(5.6), Inches(0.35),
         "Python 求解工具", sz=15, color=ORANGE, bold=True)
    py_cmd = (
        "# 批量求解文件夹\n"
        "python batch_ecxml_solver.py \\\n"
        "  ./input_folder -o ./output\n\n"
        "# FloSCRIPT 运行器 (支持功耗修改)\n"
        "python floscript_runner.py model.floxml \\\n"
        "  -o ./output --power U1_CPU=15.0\n\n"
        "# 批量求解器 (支持多种格式)\n"
        "python flotherm_batch_solver.py model.xml \\\n"
        "  -o ./output --mode auto"
    )
    _txt(s, Inches(7.0), Inches(4.7), Inches(5.6), Inches(2.0),
         py_cmd, sz=11, color=DARK_GRAY, font="Menlo")

    _footer(s)


# ========================================================================
# Slide 15: Batch simulation
# ========================================================================
def slide_batch_sim(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(s, WHITE)
    _accent_bar(s, RGBColor(0x8E, 0x44, 0xAD))
    _section_title(s, "批量求解 — batch_simulation")

    _txt(s, Inches(0.6), Inches(1.1), Inches(12), Inches(0.4),
         "batch_simulation/ 目录提供完整的批量仿真工具，通过 FloSCRIPT 宏自动加载、修改参数、求解、保存结果。",
         sz=14, color=MED_GRAY)

    # Left: workflow
    _rect(s, Inches(0.5), Inches(1.7), Inches(6.0), Inches(2.5), LIGHT_GRAY)
    _txt(s, Inches(0.7), Inches(1.8), Inches(5.6), Inches(0.35),
         "批量求解流程", sz=15, color=DARK_BLUE, bold=True)
    workflow = (
        "1. 准备 JSON 配置 (input_pack + modifications)\n"
        "2. batch_sim.py 生成 FloSCRIPT XML 脚本\n"
        "3. FlothermExecutor 调用 flotherm.exe 执行\n"
        "4. 每次求解: load → modify → solve → save_as .pack\n"
        "5. extract 命令提取结果到 JSON / CSV"
    )
    _txt(s, Inches(0.7), Inches(2.2), Inches(5.6), Inches(1.8),
         workflow, sz=12, color=DARK_GRAY)

    # Right: CLI commands
    _rect(s, Inches(6.8), Inches(1.7), Inches(6.0), Inches(2.5), LIGHT_GRAY)
    _txt(s, Inches(7.0), Inches(1.8), Inches(5.6), Inches(0.35),
         "四个子命令", sz=15, color=DARK_BLUE, bold=True)
    cli = (
        "# 生成 FloSCRIPT 脚本\n"
        "python batch_simulation/batch_sim.py \\\n"
        "  generate config.json -o ./scripts\n\n"
        "# 生成 + 执行\n"
        "python batch_simulation/batch_sim.py \\\n"
        "  run config.json --timeout 7200\n\n"
        "# 提取求解结果\n"
        "python batch_simulation/batch_sim.py \\\n"
        "  extract solved.pack -o results.json\n\n"
        "# 创建空白 FloXML 项目\n"
        "python batch_simulation/batch_sim.py \\\n"
        "  create-floxml -n MyProject -o project.xml"
    )
    _txt(s, Inches(7.0), Inches(2.2), Inches(5.6), Inches(1.8),
         cli, sz=10, color=DARK_GRAY, font="Menlo")

    # Bottom left: example configs
    _rect(s, Inches(0.5), Inches(4.4), Inches(6.0), Inches(2.8), LIGHT_GRAY)
    _txt(s, Inches(0.7), Inches(4.5), Inches(5.6), Inches(0.35),
         "example_config.json — 单次修改求解", sz=13, color=GREEN, bold=True)
    single_json = (
        '{\n'
        '  "input_pack": "model.pack",\n'
        '  "output_pack": "output.pack",\n'
        '  "modifications": [\n'
        '    {"type": "power", "component": "U1_CPU",\n'
        '     "value": 15.0},\n'
        '    {"type": "solver", "max_iterations": 500}\n'
        '  ],\n'
        '  "solve": true\n'
        '}'
    )
    _txt(s, Inches(0.7), Inches(4.9), Inches(5.6), Inches(2.0),
         single_json, sz=10, color=DARK_GRAY, font="Menlo")

    # Bottom right: param sweep
    _rect(s, Inches(6.8), Inches(4.4), Inches(6.0), Inches(2.8), LIGHT_BLUE)
    _txt(s, Inches(7.0), Inches(4.5), Inches(5.6), Inches(0.35),
         "example_param_sweep.json — 参数扫描", sz=13, color=ORANGE, bold=True)
    sweep_json = (
        '{\n'
        '  "input_pack": "model.pack",\n'
        '  "output_pack": "output_{value}.pack",\n'
        '  "parameter_sweep": {\n'
        '    "component": "U1_CPU",\n'
        '    "parameter": "power",\n'
        '    "values": [5.0, 10.0, 15.0, 20.0]\n'
        '  },\n'
        '  "solve": true\n'
        '}\n'
        "# 自动生成 4 个脚本:\n"
        "# output_5.0.pack / output_10.0.pack / ..."
    )
    _txt(s, Inches(7.0), Inches(4.9), Inches(5.6), Inches(2.0),
         sweep_json, sz=10, color=DARK_GRAY, font="Menlo")

    _footer(s)


# ========================================================================
# Slide 16: Summary
# ========================================================================
def slide_summary(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(s, WHITE)
    _accent_bar(s)
    _section_title(s, "交接总结")

    # Stats
    _rect(s, Inches(0.5), Inches(1.3), Inches(12.1), Inches(1.5), LIGHT_BLUE)
    stats = [
        ("4", "核心 Python 模块"),
        ("~5,200", "行核心代码"),
        ("6+", "JSON 配置文件"),
        ("7", "边界条件类型"),
    ]
    for i, (num, label) in enumerate(stats):
        x = Inches(0.7) + i * Inches(3.0)
        _txt(s, x, Inches(1.4), Inches(2.8), Inches(0.7),
             num, sz=36, color=ACCENT_BLUE, bold=True, align=PP_ALIGN.CENTER)
        _txt(s, x, Inches(2.0), Inches(2.8), Inches(0.4),
             label, sz=13, color=MED_GRAY, align=PP_ALIGN.CENTER)

    # Key points
    _txt(s, Inches(0.6), Inches(3.2), Inches(12), Inches(0.4),
         "核心要点", sz=20, color=DARK_BLUE, bold=True)

    points = [
        "1. 工作流: 导出 ECXML → 转换 FloXML → JSON 配置修改 → 命令行求解 → 保存 .pack + HTML 报告",
        "2. 三个配置修改脚本: solve_settings / volume_regions / boundary_conditions",
        "3. pipeline 可一键组合三个修改步骤",
        "4. 命令行求解: flotherm -b model.ecxml -z output.pack -r report.html",
        "5. 批量求解: batch_simulation/ 支持 JSON 配置驱动的参数扫描",
        "6. 所有模块都有命令行入口，支持 -h 查看帮助",
    ]
    _bullets(s, Inches(0.6), Inches(3.7), Inches(12), Inches(2.5),
             points, sz=14, color=DARK_GRAY, gap=Pt(6))

    # TODO
    _rect(s, Inches(0.5), Inches(5.6), Inches(12.1), Inches(1.3), LIGHT_GRAY)
    _txt(s, Inches(0.7), Inches(5.7), Inches(11.7), Inches(0.3),
         "后续建议", sz=14, color=DARK_BLUE, bold=True)
    todos = [
        "• 补充单元测试（目前仅 ecxml_to_floxml_converter 有测试）",
        "• 引入 logging 替代 print，便于调试",
    ]
    _bullets(s, Inches(0.7), Inches(6.0), Inches(11.7), Inches(0.8),
             todos, sz=12, color=MED_GRAY, gap=Pt(2))

    _footer(s)


# ========================================================================
# Slide 13: Thank you
# ========================================================================
def slide_thankyou(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(s, DARK_BLUE)

    _txt(s, Inches(2), Inches(2.5), Inches(9.3), Inches(1.5),
         "谢谢", sz=56, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    _txt(s, Inches(2), Inches(4.0), Inches(9.3), Inches(0.8),
         "如有疑问请参考各模块 docstring 或 README.md",
         sz=18, color=MID_BLUE, align=PP_ALIGN.CENTER)
    _txt(s, Inches(2), Inches(5.0), Inches(9.3), Inches(0.6),
         "python -m floxml_tools.<module_name> -h",
         sz=14, color=LIGHT_BLUE, align=PP_ALIGN.CENTER, font="Menlo")


# ========================================================================
def main():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    slide_cover(prs)
    slide_workflow(prs)
    slide_ecxml_converter(prs)
    slide_solve_settings(prs)
    slide_solve_config(prs)
    slide_volume_regions(prs)
    slide_volume_config(prs)
    slide_boundary_conditions(prs)
    slide_boundary_config(prs)
    slide_typical_usage(prs)
    slide_pipeline(prs)
    slide_config_summary(prs)
    slide_xsd_overview(prs)
    slide_xsd_details(prs)
    slide_cli_solve(prs)
    slide_batch_sim(prs)
    slide_summary(prs)
    slide_thankyou(prs)

    out_path = os.path.join(os.path.dirname(os.path.abspath(__file__)),
                            "floxml_tools_工作交接.pptx")
    prs.save(out_path)
    print(f"PPT saved: {out_path}")
    print(f"Total slides: {len(prs.slides)}")


if __name__ == "__main__":
    main()
